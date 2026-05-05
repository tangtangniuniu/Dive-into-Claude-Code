"""
Generate fully-editable PPTX of "Claude Code 技术与架构剖析".

Every shape is a native pptx object (text frame, table, autoshape, line, group).
Open in PowerPoint / Keynote / LibreOffice — every page can be edited.

Usage:
    python3 generate_pptx.py
Output:
    claudecode_ppt.pptx (next to this script)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree
import os

# -------- 16:9 canvas (13.333" x 7.5") --------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# -------- Color palette --------
INK        = RGBColor(0x1C, 0x25, 0x32)
INK2       = RGBColor(0x3C, 0x47, 0x57)
MUTED      = RGBColor(0x6B, 0x72, 0x80)
ACCENT     = RGBColor(0xDF, 0x7E, 0x3C)
ACCENT2    = RGBColor(0xC9, 0x74, 0x2A)
TEAL       = RGBColor(0x3D, 0x80, 0x33)
DEEP_NAVY  = RGBColor(0x1C, 0x25, 0x32)
DEEP_NAVY2 = RGBColor(0x2A, 0x36, 0x50)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG   = RGBColor(0xFA, 0xFB, 0xFD)
BORDER     = RGBColor(0xD6, 0xDA, 0xE0)
ORANGE_BG  = RGBColor(0xFD, 0xE2, 0xCC)
YELLOW_BG  = RGBColor(0xFF, 0xF1, 0xB8)
GREEN_BG   = RGBColor(0xD8, 0xEC, 0xE2)
BLUE_BG    = RGBColor(0xE1, 0xEC, 0xF5)
GRAY_BG    = RGBColor(0xEF, 0xEF, 0xEF)
ROSE_BG    = RGBColor(0xFF, 0xE4, 0xD2)
BROWN_TXT  = RGBColor(0x3A, 0x22, 0x08)
RED_WARN   = RGBColor(0xC9, 0x3A, 0x1C)


# ---- font default (use a CJK-aware safe family chain) ----
DEFAULT_FONT = "Microsoft YaHei"
FALLBACK_FONT = "Arial"


# =====================================================================
# Helpers
# =====================================================================

def set_run(run, text, *, size=14, bold=False, italic=False, color=INK,
            font=DEFAULT_FONT, mono=False):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color
    run.font.name = "Consolas" if mono else font
    # also set East-Asian font so CJK renders correctly
    rPr = run._r.get_or_add_rPr()
    # remove any prior <a:ea>
    for old in rPr.findall(qn("a:ea")):
        rPr.remove(old)
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", "Consolas" if mono else font)
    return run


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, italic=False,
             color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=DEFAULT_FONT, mono=False, line_spacing=1.2):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    set_run(run, text, size=size, bold=bold, italic=italic, color=color,
            font=font, mono=mono)
    return box


def add_multiline(slide, x, y, w, h, parts, *, align=PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.TOP, line_spacing=1.3):
    """parts = list of dicts: {'text','size','bold','color','italic','mono','newline','bullet'}.
       newline=True starts a fresh paragraph; bullet=True adds a small dot prefix.
    """
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    first = True
    for it in parts:
        if it.get("newline") and not first:
            p = tf.add_paragraph()
            p.alignment = align
            p.line_spacing = line_spacing
        first = False
        if it.get("bullet"):
            r = p.add_run()
            set_run(r, "• ", size=it.get("size", 14), bold=False,
                    color=it.get("color", INK))
        r = p.add_run()
        set_run(r, it["text"],
                size=it.get("size", 14),
                bold=it.get("bold", False),
                italic=it.get("italic", False),
                color=it.get("color", INK),
                mono=it.get("mono", False))
    return box


def set_fill(shape, color, *, transparent=False):
    if transparent:
        shape.fill.background()
        return
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color=None, width_pt=None):
    if color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = color
    if width_pt is not None:
        shape.line.width = Pt(width_pt)


def add_rect(slide, x, y, w, h, *, fill=WHITE, line=BORDER, line_pt=1.0,
             rounded=False, radius_ratio=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    rect = slide.shapes.add_shape(shape_type, x, y, w, h)
    if rounded and radius_ratio is not None:
        # adjust corner radius via XML
        try:
            rect.adjustments[0] = radius_ratio
        except Exception:
            pass
    rect.shadow.inherit = False
    set_fill(rect, fill)
    set_line(rect, line, line_pt)
    rect.text_frame.text = ""  # ensure no placeholder
    return rect


def add_pill(slide, x, y, w, h, text, *, fill=GRAY_BG, color=INK2, size=10,
             bold=True):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    pill.adjustments[0] = 0.5
    pill.shadow.inherit = False
    set_fill(pill, fill)
    set_line(pill, None)
    tf = pill.text_frame
    tf.margin_left = tf.margin_right = Emu(45000)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, text, size=size, bold=bold, color=color)
    return pill


def add_line(slide, x1, y1, x2, y2, *, color=INK, width_pt=1.5):
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
    set_line(ln, color, width_pt)
    return ln


def add_top_bar(slide, h=Inches(0.1)):
    """Orange top accent bar."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, h)
    bar.shadow.inherit = False
    set_fill(bar, ACCENT)
    set_line(bar, None)
    return bar


def add_section_pill(slide, x, y, text):
    add_text(slide, x, y, Inches(6), Inches(0.3), text,
             size=12, bold=True, color=ACCENT2)


def add_footer(slide, left_text, page_text):
    add_text(slide, Inches(0.5), Inches(7.15), Inches(8), Inches(0.3),
             left_text, size=10, color=MUTED)
    add_text(slide, Inches(11.5), Inches(7.15), Inches(1.5), Inches(0.3),
             page_text, size=10, color=MUTED, align=PP_ALIGN.RIGHT, bold=True)


def add_title(slide, x, y, text, *, size=28, color=INK, w=Inches(12), h=Inches(0.9)):
    add_text(slide, x, y, w, h, text, size=size, bold=True, color=color)


# =====================================================================
# Table helpers
# =====================================================================

def add_table(slide, x, y, w, h, headers, rows, *,
              header_fill=RGBColor(0xF6, 0xF3, 0xEE),
              header_color=BROWN_TXT,
              first_col_bold=True, body_size=11, header_size=12):
    cols = len(headers)
    rs = len(rows) + 1
    tbl_shape = slide.shapes.add_table(rs, cols, x, y, w, h)
    table = tbl_shape.table

    # header row
    for c, h_text in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        tf = cell.text_frame
        tf.margin_left = Emu(72000)
        tf.margin_right = Emu(72000)
        tf.margin_top = Emu(36000)
        tf.margin_bottom = Emu(36000)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        set_run(r, h_text, size=header_size, bold=True, color=header_color)

    # body
    for ri, row in enumerate(rows, start=1):
        for ci, cell_val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 1 else LIGHT_BG
            tf = cell.text_frame
            tf.margin_left = Emu(72000)
            tf.margin_right = Emu(72000)
            tf.margin_top = Emu(36000)
            tf.margin_bottom = Emu(36000)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT

            if isinstance(cell_val, list):
                first = True
                for seg in cell_val:
                    if seg.get("newline") and not first:
                        p = tf.add_paragraph()
                        p.alignment = PP_ALIGN.LEFT
                    first = False
                    r = p.add_run()
                    set_run(r, seg["text"],
                            size=seg.get("size", body_size),
                            bold=seg.get("bold", False),
                            italic=seg.get("italic", False),
                            color=seg.get("color", INK),
                            mono=seg.get("mono", False))
            else:
                r = p.add_run()
                set_run(r, str(cell_val), size=body_size,
                        bold=(first_col_bold and ci == 0), color=INK)

    return tbl_shape


# =====================================================================
# Slide builders
# =====================================================================

def slide_1_cover(prs):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.shadow.inherit = False
    set_fill(bg, DEEP_NAVY)
    set_line(bg, None)
    # accent radial-ish patch (using a soft orange rectangle in corner)
    patch = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(8), Inches(3.0),
                                Inches(7), Inches(7))
    patch.shadow.inherit = False
    set_fill(patch, RGBColor(0xFF, 0xD9, 0xA8))
    set_line(patch, None)
    patch.fill.transparency = 0.5  # not all viewers respect; benign
    add_top_bar(s)

    add_text(s, Inches(0.5), Inches(0.85), Inches(8), Inches(0.4),
             "DEEP-DIVE · CLAUDE CODE v2.1.88",
             size=12, bold=True, color=RGBColor(0xFF, 0xD9, 0xA8))

    add_text(s, Inches(0.5), Inches(1.4), Inches(11), Inches(1.2),
             "Claude Code", size=54, bold=True, color=WHITE)
    add_text(s, Inches(0.5), Inches(2.35), Inches(11), Inches(1.2),
             "技术与架构剖析", size=44, bold=True, color=WHITE)

    add_text(s, Inches(0.5), Inches(3.55), Inches(11), Inches(0.6),
             "A While-Loop Wrapped in 98.4% Infrastructure",
             size=18, bold=True, color=RGBColor(0xFF, 0xD9, 0xA8))

    add_text(s, Inches(0.5), Inches(4.3), Inches(9), Inches(1.6),
             "面向 v2.1.88（约 1,884 个文件 / 约 512K 行 TypeScript）的源码级架构分析。"
             "智能体循环只是一个 while；真正的工程复杂度集中在它周围的权限门控、"
             "上下文管理、工具路由与恢复逻辑等基础设施之中。",
             size=14, color=RGBColor(0xD8, 0xDD, 0xE7), line_spacing=1.55)

    badges = [
        ("1.6%", "AI 决策"),
        ("98.4%", "基础设施"),
        ("7", "安全层"),
        ("5", "压缩阶段"),
        ("54", "工具"),
        ("27", "钩子事件"),
        ("4", "扩展机制"),
        ("7", "权限模式"),
    ]
    bx = Inches(0.5)
    by = Inches(6.25)
    bw = Inches(1.45)
    bh = Inches(0.55)
    gap = Inches(0.05)
    for i, (n, k) in enumerate(badges):
        col = i % 8
        x = bx + (bw + gap) * col
        rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, by, bw, bh)
        rect.adjustments[0] = 0.18
        rect.shadow.inherit = False
        set_fill(rect, RGBColor(0x2C, 0x37, 0x4E))
        set_line(rect, RGBColor(0x4A, 0x55, 0x6B), 0.75)
        tf = rect.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r1 = p.add_run()
        set_run(r1, n, size=13, bold=True, color=ACCENT)
        r2 = p.add_run()
        set_run(r2, "  " + k, size=11, color=WHITE)

    add_text(s, Inches(11.5), Inches(7.15), Inches(1.5), Inches(0.3),
             "01 / 15", size=10, color=RGBColor(0xD8, 0xDD, 0xE7),
             align=PP_ALIGN.RIGHT, bold=True)


def slide_2_tldr(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(s)
    add_section_pill(s, Inches(0.5), Inches(0.45), "TL;DR · 一图概览")
    add_title(s, Inches(0.5), Inches(0.85),
              "智能体循环只占代码的 1.6%，真正的工程在它的\"骨架\"里",
              size=28)
    add_text(s, Inches(0.5), Inches(1.85), Inches(12.3), Inches(0.85),
             "Claude Code 不是一个\"更聪明的模型\"，而是围绕模型构建的一整套确定性基础设施："
             "权限门控、上下文管理、工具路由、恢复逻辑——这些才是从研究原型到生产可用的关键。",
             size=14, color=INK2, line_spacing=1.55)

    # stats row
    stats = [("1.6%", "AI 决策逻辑"), ("98.4%", "确定性基础设施"),
             ("~512K", "行代码 (TS)"), ("1,884", "源文件"),
             ("~200K", "上下文 token"), ("5", "压缩阶段")]
    sx = Inches(0.5); sy = Inches(2.95); sw = Inches(2.0); sh = Inches(1.1); gap = Inches(0.10)
    for i, (v, k) in enumerate(stats):
        x = sx + (sw + gap) * i
        rect = add_rect(s, x, sy, sw, sh, fill=LIGHT_BG, line=BORDER, line_pt=0.75,
                        rounded=True, radius_ratio=0.10)
        tf = rect.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); set_run(r, v, size=24, bold=True, color=ACCENT)
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); set_run(r2, k, size=10, color=MUTED)

    # two cards
    cx = Inches(0.5); cy = Inches(4.35); cw = Inches(6.15); ch = Inches(2.2)
    c1 = add_rect(s, cx, cy, cw, ch, fill=WHITE, line=BORDER, line_pt=0.75,
                  rounded=True, radius_ratio=0.04)
    bar1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, cy, Inches(0.08), ch)
    bar1.shadow.inherit = False; set_fill(bar1, ACCENT); set_line(bar1, None)
    add_text(s, cx + Inches(0.25), cy + Inches(0.15), cw - Inches(0.3), Inches(0.35),
             "核心洞察", size=11, bold=True, color=MUTED)
    add_multiline(s, cx + Inches(0.25), cy + Inches(0.55), cw - Inches(0.3), ch - Inches(0.7),
                  [{"text": "循环本身容易复制", "size": 14, "bold": True},
                   {"text": "，但 harness（钩子、分类器、压缩、隔离）才是难以重新实现的部分。",
                    "size": 14}], line_spacing=1.55)

    c2x = cx + cw + Inches(0.2)
    c2 = add_rect(s, c2x, cy, cw, ch, fill=WHITE, line=BORDER, line_pt=0.75,
                  rounded=True, radius_ratio=0.04)
    bar2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, c2x, cy, Inches(0.08), ch)
    bar2.shadow.inherit = False; set_fill(bar2, TEAL); set_line(bar2, None)
    add_text(s, c2x + Inches(0.25), cy + Inches(0.15), cw - Inches(0.3), Inches(0.35),
             "设计哲学", size=11, bold=True, color=MUTED)
    add_multiline(s, c2x + Inches(0.25), cy + Inches(0.55), cw - Inches(0.3), ch - Inches(0.7),
                  [{"text": "价值观优先于规则", "size": 14, "bold": True},
                   {"text": "：5 个人类价值观 → 13 条原则 → 落地实现，每个决策都可追溯。",
                    "size": 14}], line_spacing=1.55)

    add_footer(s, "Claude Code 技术与架构剖析", "02 / 15")


def slide_3_four_questions(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(s)
    add_section_pill(s, Inches(0.5), Inches(0.45), "CHAPTER 01 · 设计空间")
    add_title(s, Inches(0.5), Inches(0.85),
              "每个编码智能体都必须回答的四个问题", size=28)

    headers = ["设计问题", "Claude Code 的答案", "替代方案"]
    rows = [
        [
            [{"text": "推理放在哪里？", "bold": True}],
            [{"text": "模型推理；harness 强制执行", "bold": False},
             {"text": "≈1.6% AI · 98.4% 基础设施", "newline": True,
              "color": ACCENT, "bold": True, "size": 11}],
            "LangGraph：显式状态图\nDevin：多步规划器",
        ],
        [
            [{"text": "有多少个执行引擎？", "bold": True}],
            [{"text": "单一 ", "size": 12},
             {"text": "queryLoop", "mono": True, "size": 12, "color": ACCENT2},
             {"text": "，所有入口（CLI / SDK / IDE）共用", "size": 12},
             {"text": "One harness for all", "newline": True,
              "color": RGBColor(0x23, 0x4A, 0x6F), "bold": True, "size": 11}],
            "每个入口一套独立引擎",
        ],
        [
            [{"text": "默认安全姿态？", "bold": True}],
            [{"text": "拒绝优先：拒绝 > 询问 > 允许；最严格规则胜出", "size": 12},
             {"text": "Deny-first", "newline": True,
              "color": RGBColor(0x5A, 0x44, 0x00), "bold": True, "size": 11}],
            "容器隔离（SWE-Agent）\nGit 回滚（Aider）",
        ],
        [
            [{"text": "最根本的资源约束？", "bold": True}],
            [{"text": "上下文窗口（~200K，部分 1M）；每次模型调用前过 5 层压缩", "size": 12},
             {"text": "5-stage shaping", "newline": True,
              "color": TEAL, "bold": True, "size": 11}],
            "计算预算、显式草稿本",
        ],
    ]
    add_table(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(4.6),
              headers, rows, body_size=12)
    # column widths
    tbl = s.shapes[-1].table
    tbl.columns[0].width = Inches(2.6)
    tbl.columns[1].width = Inches(5.5)
    tbl.columns[2].width = Inches(4.2)
    add_footer(s, "四个设计问题 · From Paper § 2", "03 / 15")


def slide_4_main_structure(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(s)
    add_section_pill(s, Inches(0.5), Inches(0.45), "CHAPTER 02 · 系统总览")
    add_title(s, Inches(0.5), Inches(0.85),
              "高层主结构 · 7 个组件构成的请求—执行—持久化闭环", size=26)
    add_text(s, Inches(0.5), Inches(1.85), Inches(12.3), Inches(0.7),
             "User → Interfaces → Agent Loop（核心 while），由 Permission System 门控、调度 "
             "Tools 在 Execution Environment 中运行，状态写入 State & Persistence。",
             size=13, color=INK2, line_spacing=1.5)

    # 7 component boxes laid out per main_structure.svg
    # Row layout:
    #  Row1 (top middle): Permission System
    #  Row2 (mid):  User -> Interfaces -> Agent Loop -> Tools -> Execution Environment
    #  Row3 (bottom): State & Persistence (under Agent Loop)

    # Permission System (top center)
    add_rect(s, Inches(5.6), Inches(2.5), Inches(2.6), Inches(0.95),
             fill=YELLOW_BG, line=RGBColor(0xC9, 0xA3, 0x00), line_pt=2.0,
             rounded=True, radius_ratio=0.20)
    add_text(s, Inches(5.6), Inches(2.55), Inches(2.6), Inches(0.4),
             "Permission System", size=15, bold=True, color=INK,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(5.6), Inches(2.95), Inches(2.6), Inches(0.4),
             "+ auto classifier", size=11, color=BROWN_TXT, italic=True,
             align=PP_ALIGN.CENTER)

    # User
    add_rect(s, Inches(0.6), Inches(4.5), Inches(1.4), Inches(1.0),
             fill=GRAY_BG, line=INK, line_pt=2.0,
             rounded=True, radius_ratio=0.10)
    add_text(s, Inches(0.6), Inches(4.85), Inches(1.4), Inches(0.4),
             "User", size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)

    # Interfaces
    add_rect(s, Inches(2.4), Inches(4.4), Inches(2.0), Inches(1.2),
             fill=ORANGE_BG, line=ACCENT2, line_pt=2.0,
             rounded=True, radius_ratio=0.20)
    add_text(s, Inches(2.4), Inches(4.85), Inches(2.0), Inches(0.4),
             "Interfaces", size=15, bold=True, align=PP_ALIGN.CENTER)

    # Agent Loop
    add_rect(s, Inches(4.9), Inches(4.4), Inches(2.4), Inches(1.4),
             fill=ACCENT, line=RGBColor(0xA6, 0x4F, 0x17), line_pt=2.0,
             rounded=True, radius_ratio=0.18)
    add_text(s, Inches(4.9), Inches(4.65), Inches(2.4), Inches(0.4),
             "Agent", size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(4.9), Inches(5.05), Inches(2.4), Inches(0.4),
             "Loop", size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)

    # Tools
    add_rect(s, Inches(7.8), Inches(4.5), Inches(1.7), Inches(1.0),
             fill=BLUE_BG, line=RGBColor(0x3A, 0x72, 0xC4), line_pt=2.0,
             rounded=True, radius_ratio=0.20)
    add_text(s, Inches(7.8), Inches(4.85), Inches(1.7), Inches(0.4),
             "Tools", size=16, bold=True, align=PP_ALIGN.CENTER)

    # Execution Environment
    add_rect(s, Inches(10.0), Inches(4.4), Inches(2.6), Inches(1.4),
             fill=LIGHT_BG, line=MUTED, line_pt=2.0,
             rounded=True, radius_ratio=0.16)
    add_text(s, Inches(10.0), Inches(4.6), Inches(2.6), Inches(0.4),
             "Execution", size=13, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(10.0), Inches(4.95), Inches(2.6), Inches(0.4),
             "Environment", size=13, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(10.0), Inches(5.4), Inches(2.6), Inches(0.3),
             "Files / Shell / Web / MCP", size=10, color=MUTED,
             italic=True, align=PP_ALIGN.CENTER)

    # State & Persistence
    add_rect(s, Inches(4.9), Inches(6.05), Inches(2.4), Inches(1.0),
             fill=GREEN_BG, line=RGBColor(0x3D, 0x80, 0x33), line_pt=2.0,
             rounded=True, radius_ratio=0.20)
    add_text(s, Inches(4.9), Inches(6.15), Inches(2.4), Inches(0.4),
             "State &", size=14, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(4.9), Inches(6.55), Inches(2.4), Inches(0.4),
             "Persistence", size=14, bold=True, align=PP_ALIGN.CENTER)

    # Connectors
    arr = lambda x1, y1, x2, y2: add_line(s, x1, y1, x2, y2, color=INK, width_pt=1.6)
    arr(Inches(2.0), Inches(5.0), Inches(2.4), Inches(5.0))   # User -> IF
    arr(Inches(4.4), Inches(5.0), Inches(4.9), Inches(5.0))   # IF -> Agent
    arr(Inches(7.3), Inches(5.0), Inches(7.8), Inches(5.0))   # Agent -> Tools
    arr(Inches(9.5), Inches(5.0), Inches(10.0), Inches(5.0))  # Tools -> Env
    # Agent -> Permission (up)
    arr(Inches(6.1), Inches(4.4), Inches(6.1), Inches(3.45))
    # Permission -> Tools (curve approximated by line)
    arr(Inches(7.3), Inches(3.45), Inches(8.6), Inches(4.5))
    # Agent -> State (down) and back
    arr(Inches(5.7), Inches(5.8), Inches(5.7), Inches(6.05))
    arr(Inches(6.5), Inches(6.05), Inches(6.5), Inches(5.8))

    # labels (italic, light)
    add_text(s, Inches(2.0), Inches(4.4), Inches(0.6), Inches(0.3),
             "Prompt", size=9, italic=True, color=INK2, align=PP_ALIGN.CENTER)
    add_text(s, Inches(4.4), Inches(4.4), Inches(0.6), Inches(0.3),
             "Request", size=9, italic=True, color=INK2, align=PP_ALIGN.CENTER)
    add_text(s, Inches(7.3), Inches(4.4), Inches(0.6), Inches(0.3),
             "Tool call", size=9, italic=True, color=INK2, align=PP_ALIGN.CENTER)
    add_text(s, Inches(6.0), Inches(3.5), Inches(2.0), Inches(0.3),
             "Propose", size=9, italic=True, color=INK2, align=PP_ALIGN.CENTER)
    add_text(s, Inches(7.4), Inches(3.7), Inches(1.6), Inches(0.3),
             "Allow/Ask/Deny", size=9, italic=True, color=INK2)
    add_text(s, Inches(5.7), Inches(5.85), Inches(0.7), Inches(0.3),
             "Persist", size=9, italic=True, color=INK2)
    add_text(s, Inches(6.55), Inches(5.85), Inches(0.6), Inches(0.3),
             "Load", size=9, italic=True, color=INK2)

    add_footer(s, "主结构 · 7 个组件", "04 / 15")


def slide_5_layered(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(s)
    add_section_pill(s, Inches(0.5), Inches(0.45), "CHAPTER 02 · 系统总览")
    add_title(s, Inches(0.5), Inches(0.85),
              "5 层子系统分解 · 21 个子系统", size=26)
    add_text(s, Inches(0.5), Inches(1.85), Inches(12.3), Inches(0.7),
             "Surface（入口/UI）· Core（智能体循环 + 压缩）· Safety/Action（权限 + 钩子 + 工具）· "
             "Backend（执行后端）；底部 State 贯穿全部层。",
             size=13, color=INK2, line_spacing=1.5)

    # Top row: 4 layer panels
    row_y = Inches(2.6); row_h = Inches(3.5)
    # Surface
    add_rect(s, Inches(0.5), row_y, Inches(2.6), row_h,
             fill=BLUE_BG, line=RGBColor(0x88, 0xA4, 0xBD), line_pt=1.4,
             rounded=True, radius_ratio=0.04)
    add_text(s, Inches(0.6), row_y + Inches(0.1), Inches(2.4), Inches(0.4),
             "▶ Surface Layer", size=13, bold=True)

    # 4 sources
    src_names = ["Interactive CLI", "IDE/Desktop/Browser", "Headless CLI", "Agent SDK"]
    for i, name in enumerate(src_names):
        sy = row_y + Inches(0.55) + Inches(0.62) * i
        add_rect(s, Inches(0.65), sy, Inches(1.5), Inches(0.5),
                 fill=WHITE, line=RGBColor(0x5D, 0x7E, 0x9D), line_pt=0.75,
                 rounded=True, radius_ratio=0.18)
        add_text(s, Inches(0.65), sy + Inches(0.1), Inches(1.5), Inches(0.3),
                 name, size=10, bold=True, align=PP_ALIGN.CENTER)
    # UI/Render hub
    add_rect(s, Inches(2.3), row_y + Inches(1.0), Inches(0.7), Inches(2.0),
             fill=WHITE, line=RGBColor(0x34, 0x51, 0x6F), line_pt=1.0,
             rounded=True, radius_ratio=0.15)
    add_text(s, Inches(2.3), row_y + Inches(1.4), Inches(0.7), Inches(0.4),
             "UI /", size=11, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2.3), row_y + Inches(1.7), Inches(0.7), Inches(0.4),
             "Render", size=11, bold=True, align=PP_ALIGN.CENTER)

    # Core
    add_rect(s, Inches(3.25), row_y, Inches(2.0), row_h,
             fill=RGBColor(0xDD, 0xE8, 0xF3),
             line=RGBColor(0x7A, 0x99, 0xB8), line_pt=1.4,
             rounded=True, radius_ratio=0.06)
    add_text(s, Inches(3.35), row_y + Inches(0.1), Inches(2.0), Inches(0.4),
             "▶ Core Layer", size=13, bold=True)
    # Agent Loop
    add_rect(s, Inches(3.4), row_y + Inches(0.7), Inches(1.7), Inches(0.9),
             fill=WHITE, line=RGBColor(0x34, 0x51, 0x6F), line_pt=1.2,
             rounded=True, radius_ratio=0.12)
    add_text(s, Inches(3.4), row_y + Inches(0.95), Inches(1.7), Inches(0.4),
             "Agent Loop", size=14, bold=True, align=PP_ALIGN.CENTER)
    # Compaction
    add_rect(s, Inches(3.4), row_y + Inches(2.4), Inches(1.7), Inches(0.85),
             fill=WHITE, line=RGBColor(0x34, 0x51, 0x6F), line_pt=1.2,
             rounded=True, radius_ratio=0.12)
    add_text(s, Inches(3.4), row_y + Inches(2.5), Inches(1.7), Inches(0.4),
             "Compaction", size=11, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(3.4), row_y + Inches(2.85), Inches(1.7), Inches(0.4),
             "Pipeline", size=11, bold=True, align=PP_ALIGN.CENTER)
    # arrows
    add_line(s, Inches(3.9), row_y + Inches(1.6), Inches(3.9), row_y + Inches(2.4),
             color=INK, width_pt=1.2)
    add_line(s, Inches(4.6), row_y + Inches(2.4), Inches(4.6), row_y + Inches(1.6),
             color=INK, width_pt=1.2)

    # Safety/Action
    add_rect(s, Inches(5.4), row_y, Inches(4.9), row_h,
             fill=GREEN_BG, line=RGBColor(0x69, 0xA2, 0x88), line_pt=1.4,
             rounded=True, radius_ratio=0.04)
    add_text(s, Inches(5.5), row_y + Inches(0.1), Inches(4.0), Inches(0.4),
             "▶ Safety / Action Layer", size=13, bold=True)

    # Permission box
    add_rect(s, Inches(5.6), row_y + Inches(0.65), Inches(1.9), Inches(0.85),
             fill=ORANGE_BG, line=RGBColor(0xC9, 0x74, 0x2A), line_pt=1.0,
             rounded=True, radius_ratio=0.12)
    add_text(s, Inches(5.6), row_y + Inches(0.8), Inches(1.9), Inches(0.3),
             "Permission System", size=11, bold=True, color=BROWN_TXT,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(5.6), row_y + Inches(1.1), Inches(1.9), Inches(0.3),
             "+ auto classifier", size=10, color=BROWN_TXT, italic=True,
             align=PP_ALIGN.CENTER)
    # Hook Pipeline
    add_rect(s, Inches(5.6), row_y + Inches(1.95), Inches(1.9), Inches(0.7),
             fill=YELLOW_BG, line=RGBColor(0xCA, 0xA8, 0x00), line_pt=1.0,
             rounded=True, radius_ratio=0.14)
    add_text(s, Inches(5.6), row_y + Inches(2.1), Inches(1.9), Inches(0.4),
             "Hook Pipeline", size=12, bold=True, color=BROWN_TXT,
             align=PP_ALIGN.CENTER)

    # Right: 4 tools/extensions
    rr_x = Inches(7.8); rr_w = Inches(2.4)
    rights = [("Extensibility", "(plugins & skills)", row_y + Inches(0.55)),
              ("Built-in Tools", "", row_y + Inches(1.25)),
              ("MCP Tools", "", row_y + Inches(1.95)),
              ("Subagent spawning", "", row_y + Inches(2.65))]
    for name, sub, y in rights:
        add_rect(s, rr_x, y, rr_w, Inches(0.55),
                 fill=WHITE, line=RGBColor(0x3D, 0x7C, 0x5A), line_pt=0.9,
                 rounded=True, radius_ratio=0.16)
        if sub:
            add_text(s, rr_x, y + Inches(0.05), rr_w, Inches(0.3),
                     name, size=11, bold=True, align=PP_ALIGN.CENTER)
            add_text(s, rr_x, y + Inches(0.28), rr_w, Inches(0.3),
                     sub, size=8, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
        else:
            add_text(s, rr_x, y + Inches(0.13), rr_w, Inches(0.3),
                     name, size=11, bold=True, align=PP_ALIGN.CENTER)

    # Backend
    add_rect(s, Inches(10.4), row_y, Inches(2.4), row_h,
             fill=RGBColor(0xE2, 0xE6, 0xEB), line=MUTED, line_pt=1.4,
             rounded=True, radius_ratio=0.04)
    add_text(s, Inches(10.5), row_y + Inches(0.1), Inches(2.2), Inches(0.4),
             "▶ Backend Layer", size=13, bold=True)
    add_rect(s, Inches(10.55), row_y + Inches(0.65), Inches(2.1), Inches(0.55),
             fill=WHITE, line=MUTED, line_pt=1.0, rounded=True, radius_ratio=0.16)
    add_text(s, Inches(10.55), row_y + Inches(0.78), Inches(2.1), Inches(0.4),
             "Execution Backends", size=11, bold=True, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(10.55), row_y + Inches(1.55), Inches(2.1), Inches(0.55),
             fill=WHITE, line=MUTED, line_pt=1.0, rounded=True, radius_ratio=0.16)
    add_text(s, Inches(10.55), row_y + Inches(1.68), Inches(2.1), Inches(0.4),
             "External Resources", size=11, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(10.55), row_y + Inches(2.3), Inches(2.1), Inches(0.3),
             "local / cloud / remote", size=10, italic=True, color=MUTED,
             align=PP_ALIGN.CENTER)

    # Bottom State Layer (full width)
    state_y = Inches(6.25)
    add_rect(s, Inches(0.5), state_y, Inches(12.3), Inches(0.8),
             fill=RGBColor(0xEF, 0xE2, 0xC5),
             line=RGBColor(0xA8, 0x88, 0x4A), line_pt=1.4,
             rounded=True, radius_ratio=0.10)
    add_text(s, Inches(0.65), state_y + Inches(0.05), Inches(2.0), Inches(0.3),
             "▶ State Layer", size=12, bold=True, color=BROWN_TXT)
    state_items = [
        "Context Assembly", "Runtime State", "Session Persistence",
        "CLAUDE.md + memory", "Sidechain Transcriptions"
    ]
    sx = Inches(0.7); sw = Inches(2.4); sy2 = state_y + Inches(0.4); gp = Inches(0.05)
    for i, name in enumerate(state_items):
        x = sx + (sw + gp) * i
        add_rect(s, x, sy2, sw, Inches(0.36),
                 fill=RGBColor(0xFB, 0xF3, 0xDF),
                 line=RGBColor(0x9C, 0x6A, 0x18), line_pt=0.6,
                 rounded=True, radius_ratio=0.20)
        add_text(s, x, sy2 + Inches(0.05), sw, Inches(0.3),
                 name, size=10, bold=True, color=BROWN_TXT,
                 align=PP_ALIGN.CENTER)

    add_footer(s, "分层架构 · 5 层", "05 / 15")


def slide_6_pipeline(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(s)
    add_section_pill(s, Inches(0.5), Inches(0.45), "CHAPTER 03 · 智能体循环")
    add_title(s, Inches(0.5), Inches(0.85),
              "每一轮的 9 步流水线", size=28)
    add_text(s, Inches(0.5), Inches(1.85), Inches(12.3), Inches(0.55),
             "queryLoop 是一个 AsyncGenerator，按下列 9 步推进。"
             "停止条件：无工具调用 · 最大轮次 · 上下文溢出 · 钩子干预 · 显式中止。",
             size=13, color=INK2, line_spacing=1.5)

    # 9 step pipeline
    steps = [
        ("1", "设置解析"), ("2", "状态初始化"), ("3", "上下文组装"),
        ("4", "5 阶段\n预模型整形"), ("5", "模型调用"), ("6", "工具分派"),
        ("7", "权限门控"), ("8", "工具执行"), ("9", "停止条件检查"),
    ]
    px = Inches(0.5); py = Inches(2.6); ph = Inches(1.4)
    pw = Inches(1.32); gp = Inches(0.06)
    for i, (n, name) in enumerate(steps):
        x = px + (pw + gp) * i
        rect = add_rect(s, x, py, pw, ph,
                        fill=RGBColor(0xFF, 0xE6, 0xC8),
                        line=RGBColor(0xC9, 0x74, 0x2A), line_pt=1.0,
                        rounded=True, radius_ratio=0.10)
        # circle number
        c = s.shapes.add_shape(MSO_SHAPE.OVAL,
                               x + Inches(0.45), py + Inches(0.15),
                               Inches(0.42), Inches(0.42))
        c.shadow.inherit = False
        set_fill(c, ACCENT2); set_line(c, None)
        tf = c.text_frame; tf.margin_left = tf.margin_right = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); set_run(r, n, size=12, bold=True, color=WHITE)
        # name
        add_text(s, x, py + Inches(0.65), pw, Inches(0.7),
                 name, size=11, bold=True, color=BROWN_TXT,
                 align=PP_ALIGN.CENTER, line_spacing=1.2)

    # 3 cards
    cy = Inches(4.4); cw = Inches(4.0); ch = Inches(2.2)
    cards = [
        ("两条执行路径",
         [("StreamingToolExecutor", "：边流入边执行（低延迟）"),
          ("runTools", "：分类为并发安全或互斥")]),
        ("故障恢复",
         [("最大输出 token 升级", "（≤3 次重试）"),
          ("反应式压缩", "（每轮一次）"),
          ("流式后备", " + 后备模型切换")]),
        ("5 个停止条件",
         [("无工具调用", " / 最大轮次"),
          ("上下文溢出", " / 钩子干预"),
          ("显式中止", "")]),
    ]
    for i, (title, items) in enumerate(cards):
        x = Inches(0.5) + (cw + Inches(0.15)) * i
        rect = add_rect(s, x, cy, cw, ch, fill=WHITE, line=BORDER, line_pt=0.75,
                        rounded=True, radius_ratio=0.04)
        add_text(s, x + Inches(0.2), cy + Inches(0.18), cw - Inches(0.4), Inches(0.32),
                 title, size=11, bold=True, color=MUTED)
        # bullets
        parts = []
        for j, (a, b) in enumerate(items):
            parts.append({"text": "• ", "size": 13, "color": ACCENT,
                          "newline": (j > 0)})
            parts.append({"text": a, "size": 13, "bold": True})
            parts.append({"text": b, "size": 13})
        add_multiline(s, x + Inches(0.2), cy + Inches(0.55), cw - Inches(0.4), ch - Inches(0.7),
                      parts, line_spacing=1.55)

    add_footer(s, "9 步轮次管道", "06 / 15")


def slide_7_compaction(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(s)
    add_section_pill(s, Inches(0.5), Inches(0.45), "CHAPTER 04 · 上下文管理")
    add_title(s, Inches(0.5), Inches(0.85),
              "5 阶段预模型整形 · 渐进式惰性降级", size=26)
    add_text(s, Inches(0.5), Inches(1.85), Inches(12.3), Inches(0.55),
             "每次模型调用前按顺序执行，按开销从低到高。前面阶段无法满足时才升级到下一级。",
             size=13, color=INK2, line_spacing=1.5)

    headers = ["阶段", "策略", "触发条件", "开销"]
    rows = [
        [[{"text": "① 预算削减", "bold": True}], "每条消息大小上限", "始终启用",
         [{"text": "极低", "color": TEAL, "bold": True}]],
        [[{"text": "② 裁剪 (Snip)", "bold": True}], "裁剪较旧的历史",
         [{"text": "特性开关 ", "size": 11},
          {"text": "HISTORY_SNIP", "mono": True, "size": 11, "color": ACCENT2}],
         [{"text": "低", "color": TEAL, "bold": True}]],
        [[{"text": "③ 微压缩", "bold": True}], "缓存感知的细粒度压缩",
         "始终启用（基于时间）",
         [{"text": "中", "color": RGBColor(0x5A, 0x44, 0x00), "bold": True}]],
        [[{"text": "④ 上下文折叠", "bold": True}], "读取时虚拟投影（非破坏性）",
         [{"text": "特性开关 ", "size": 11},
          {"text": "CONTEXT_COLLAPSE", "mono": True, "size": 11, "color": ACCENT2}],
         [{"text": "较高", "color": ACCENT2, "bold": True}]],
        [[{"text": "⑤ 自动压缩", "bold": True}],
         "完整模型生成的摘要", "其他阶段都失败时（最后手段）",
         [{"text": "高", "color": RED_WARN, "bold": True}]],
    ]
    add_table(s, Inches(0.5), Inches(2.55), Inches(12.3), Inches(3.1),
              headers, rows, body_size=12)
    tbl = s.shapes[-1].table
    tbl.columns[0].width = Inches(2.0)
    tbl.columns[1].width = Inches(4.4)
    tbl.columns[2].width = Inches(3.9)
    tbl.columns[3].width = Inches(2.0)

    # design note
    cy = Inches(5.85); ch = Inches(1.1)
    add_rect(s, Inches(0.5), cy, Inches(12.3), ch, fill=WHITE,
             line=BORDER, line_pt=0.75, rounded=True, radius_ratio=0.06)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.5), cy, Inches(0.08), ch)
    bar.shadow.inherit = False; set_fill(bar, ACCENT); set_line(bar, None)
    add_text(s, Inches(0.7), cy + Inches(0.12), Inches(11.5), Inches(0.3),
             "设计要点", size=11, bold=True, color=MUTED)
    add_multiline(s, Inches(0.7), cy + Inches(0.45), Inches(11.5), ch - Inches(0.55),
                  [{"text": "从第一天就为上下文稀缺而设计。", "size": 13, "bold": True},
                   {"text": "渐进式优于一次性截断，缓存感知优于盲目摘要——每一层都尽量保留语义保真度，把破坏性最高的\"自动压缩\"留到最后。",
                    "size": 13}], line_spacing=1.5)

    add_footer(s, "上下文管理 · 5 阶段压缩", "07 / 15")


def slide_8_defense(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(s)
    add_section_pill(s, Inches(0.5), Inches(0.45), "CHAPTER 05 · 安全与权限")
    add_title(s, Inches(0.5), Inches(0.85),
              "深度防御 · 7 个独立安全层", size=28)
    add_text(s, Inches(0.5), Inches(1.85), Inches(12.3), Inches(0.55),
             "请求必须通过所有适用层；任一层都可拦下。\"拒绝优先\"在每一层都生效。",
             size=13, color=INK2, line_spacing=1.5)

    # 7 layers
    layers = [
        ("工具预过滤", "把全局拒绝的工具从模型可见清单中剔除"),
        ("拒绝优先规则评估", "拒绝始终覆盖允许，即使允许更具体"),
        ("权限模式约束", "活跃模式决定基线处理（plan / default / acceptEdits / auto / dontAsk / bypass）"),
        ("Auto 模式 ML 分类器", "独立 LLM 调用：快速过滤 + 思维链评估"),
        ("Shell 沙箱", "对 shell 命令实施文件系统 + 网络隔离"),
        ("恢复会话权限永不自动恢复", "权限绝不跨会话保留"),
        ("钩子拦截", "PreToolUse 钩子可修改或阻止操作"),
    ]
    fills = [RGBColor(0xFF, 0xE4, 0xD2), RGBColor(0xFF, 0xD9, 0xBE),
             RGBColor(0xFF, 0xCE, 0xAA), RGBColor(0xFF, 0xC2, 0x96),
             RGBColor(0xFF, 0xB7, 0x81), RGBColor(0xF6, 0xA8, 0x6A),
             RGBColor(0xEC, 0x9A, 0x55)]
    ly = Inches(2.55); lh = Inches(0.55); lgap = Inches(0.05)
    for i, ((title, desc), fill) in enumerate(zip(layers, fills)):
        y = ly + (lh + lgap) * i
        rect = add_rect(s, Inches(0.5), y, Inches(12.3), lh, fill=fill,
                        line=None, line_pt=0, rounded=True, radius_ratio=0.20)
        # number circle
        c = s.shapes.add_shape(MSO_SHAPE.OVAL,
                               Inches(0.7), y + Inches(0.07),
                               Inches(0.42), Inches(0.42))
        c.shadow.inherit = False
        set_fill(c, ACCENT); set_line(c, None)
        tf = c.text_frame; tf.margin_left = tf.margin_right = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); set_run(r, str(i + 1), size=14, bold=True, color=WHITE)
        # texts
        add_multiline(s, Inches(1.3), y + Inches(0.13), Inches(10.9), Inches(0.32),
                      [{"text": title, "size": 13, "bold": True, "color": INK},
                       {"text": "    ", "size": 13},
                       {"text": desc, "size": 12, "color": BROWN_TXT}])

    add_footer(s, "七层安全防御", "08 / 15")


def slide_9_modes(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(s)
    add_section_pill(s, Inches(0.5), Inches(0.45), "CHAPTER 05 · 安全与权限")
    add_title(s, Inches(0.5), Inches(0.85),
              "渐进式信任光谱 · 7 个权限模式", size=26)

    headers = ["模式", "行为", "信任级别"]
    rows = [
        [[{"text": "plan", "mono": True, "color": ACCENT2}], "用户在执行前批准所有计划",
         [{"text": "最低", "color": RGBColor(0x7A, 0x2C, 0x08), "bold": True}]],
        [[{"text": "default", "mono": True, "color": ACCENT2}], "标准交互批准",
         [{"text": "低", "color": RGBColor(0x7A, 0x2C, 0x08), "bold": True}]],
        [[{"text": "acceptEdits", "mono": True, "color": ACCENT2}],
         "文件编辑 + 文件系统 shell 自动批准",
         [{"text": "中", "color": RGBColor(0x5A, 0x44, 0x00), "bold": True}]],
        [[{"text": "auto", "mono": True, "color": ACCENT2}],
         "ML 分类器评估工具安全性",
         [{"text": "高", "color": RGBColor(0x5A, 0x44, 0x00), "bold": True}]],
        [[{"text": "dontAsk", "mono": True, "color": ACCENT2}],
         "无提示，仍强制执行拒绝规则",
         [{"text": "较高", "color": TEAL, "bold": True}]],
        [[{"text": "bypassPermissions", "mono": True, "color": ACCENT2}],
         "跳过大多数提示，仍保留关键检查",
         [{"text": "最高", "color": TEAL, "bold": True}]],
        [[{"text": "bubble", "mono": True, "color": ACCENT2}],
         "内部：子智能体向父级上报",
         [{"text": "特殊", "color": INK2, "bold": True}]],
    ]
    add_table(s, Inches(0.5), Inches(1.85), Inches(12.3), Inches(3.4),
              headers, rows, body_size=11.5)
    tbl = s.shapes[-1].table
    tbl.columns[0].width = Inches(2.4)
    tbl.columns[1].width = Inches(7.6)
    tbl.columns[2].width = Inches(2.3)

    # 2 cards
    cy = Inches(5.45); ch = Inches(1.6); cw = Inches(6.05)
    add_rect(s, Inches(0.5), cy, cw, ch, fill=WHITE, line=BORDER, line_pt=0.75,
             rounded=True, radius_ratio=0.05)
    add_text(s, Inches(0.7), cy + Inches(0.15), cw - Inches(0.4), Inches(0.3),
             "授权管道（4 阶段）", size=11, bold=True, color=MUTED)
    add_multiline(s, Inches(0.7), cy + Inches(0.5), cw - Inches(0.4), Inches(0.5),
                  [{"text": "预过滤", "bold": True, "size": 12},
                   {"text": " → ", "size": 12},
                   {"text": "PreToolUse 钩子", "bold": True, "size": 12},
                   {"text": " → ", "size": 12},
                   {"text": "规则评估（拒绝优先）", "bold": True, "size": 12},
                   {"text": " → ", "size": 12},
                   {"text": "权限处理程序", "bold": True, "size": 12}])
    # pills
    pill_x = Inches(0.7); pill_y = cy + Inches(1.0); pill_w = Inches(1.3); gp = Inches(0.08)
    for i, name in enumerate(["协调器", "swarm worker", "推测式分类器", "交互式"]):
        add_pill(s, pill_x + (pill_w + gp) * i, pill_y, pill_w, Inches(0.32),
                 name, fill=BLUE_BG, color=RGBColor(0x23, 0x4A, 0x6F), size=10)

    cx2 = Inches(0.5) + cw + Inches(0.2)
    add_rect(s, cx2, cy, cw, ch, fill=WHITE, line=BORDER, line_pt=0.75,
             rounded=True, radius_ratio=0.05)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx2, cy, Inches(0.08), ch)
    bar.shadow.inherit = False; set_fill(bar, RED_WARN); set_line(bar, None)
    add_text(s, cx2 + Inches(0.25), cy + Inches(0.15), cw - Inches(0.4), Inches(0.3),
             "⚠ 共享故障模式", size=11, bold=True, color=RED_WARN)
    add_multiline(s, cx2 + Inches(0.25), cy + Inches(0.5), cw - Inches(0.4), ch - Inches(0.6),
                  [{"text": "子命令解析饥饿 REPL → 一旦超过 50 个子命令，整段安全分析会被跳过；扩展会在",
                    "size": 12},
                   {"text": "信任对话框前", "size": 12, "bold": True},
                   {"text": "就执行（4 个 CVE 来源）。", "size": 12}],
                  line_spacing=1.5)

    add_footer(s, "7 个权限模式 · 渐进信任光谱", "09 / 15")


def slide_10_extensibility(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(s)
    add_section_pill(s, Inches(0.5), Inches(0.45), "CHAPTER 06 · 可扩展性")
    add_title(s, Inches(0.5), Inches(0.85),
              "4 个扩展机制 · 上下文成本递增", size=28)

    headers = ["机制", "上下文成本", "关键能力"]
    rows = [
        [[{"text": "Hooks", "bold": True}],
         [{"text": "零", "color": TEAL, "bold": True}],
         "27 个事件 · 4 种执行类型（shell / LLM / webhook / subagent 验证器）"],
        [[{"text": "Skills", "bold": True}],
         [{"text": "低", "color": RGBColor(0x5A, 0x44, 0x00), "bold": True}],
         "SKILL.md（15+ YAML frontmatter 字段）；通过 SkillTool 元工具注入当前上下文"],
        [[{"text": "Plugins", "bold": True}],
         [{"text": "中", "color": ACCENT2, "bold": True}],
         "10 种组件类型：命令 / 智能体 / skills / hooks / MCP / LSP / 样式…"],
        [[{"text": "MCP 服务器", "bold": True}],
         [{"text": "高", "color": RED_WARN, "bold": True}],
         "外部工具 · 7 种传输类型（stdio / SSE / HTTP / WebSocket / SDK / IDE）"],
    ]
    add_table(s, Inches(0.5), Inches(1.85), Inches(12.3), Inches(2.3),
              headers, rows, body_size=12)
    tbl = s.shapes[-1].table
    tbl.columns[0].width = Inches(2.0)
    tbl.columns[1].width = Inches(2.4)
    tbl.columns[2].width = Inches(7.9)

    add_text(s, Inches(0.5), Inches(4.35), Inches(12.3), Inches(0.4),
             "智能体循环中的 3 个注入点", size=18, bold=True)

    # 3 injection points
    iy = Inches(4.85); ih = Inches(2.0); iw = Inches(4.0); gp = Inches(0.15)
    points = [
        ("① assemble()", "模型看到什么", "CLAUDE.md、skill 描述、MCP 资源、钩子注入的上下文。"),
        ("② model()", "模型能触及什么", "内置工具、MCP 工具、SkillTool、AgentTool。"),
        ("③ execute()", "操作是否/如何运行", "权限规则、PreToolUse / PostToolUse / Stop 钩子。"),
    ]
    for i, (title, sub, body) in enumerate(points):
        x = Inches(0.5) + (iw + gp) * i
        rect = add_rect(s, x, iy, iw, ih,
                        fill=YELLOW_BG, line=RGBColor(0xCA, 0xA8, 0x00), line_pt=1.0,
                        rounded=True, radius_ratio=0.06)
        add_text(s, x + Inches(0.25), iy + Inches(0.2), iw - Inches(0.5), Inches(0.4),
                 title, size=16, bold=True, color=RGBColor(0x5A, 0x44, 0x00))
        add_text(s, x + Inches(0.25), iy + Inches(0.65), iw - Inches(0.5), Inches(0.4),
                 sub, size=12, bold=True, italic=True, color=BROWN_TXT)
        add_text(s, x + Inches(0.25), iy + Inches(1.05), iw - Inches(0.5), Inches(0.85),
                 body, size=11, color=BROWN_TXT, line_spacing=1.5)

    add_footer(s, "4 个扩展机制 · 3 个注入点", "10 / 15")


def slide_11_context_memory(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(s)
    add_section_pill(s, Inches(0.5), Inches(0.45), "CHAPTER 07 · 上下文与记忆")
    add_title(s, Inches(0.5), Inches(0.85),
              "9 个有序上下文来源 + 4 级 CLAUDE.md 层级", size=26)

    # left: 9 sources
    add_text(s, Inches(0.5), Inches(2.0), Inches(5.5), Inches(0.4),
             "9 个有序来源", size=18, bold=True)
    sources = ["系统提示", "环境信息", "CLAUDE.md 层级", "按路径限定的规则",
               "自动记忆", "工具元数据", "对话历史", "工具结果", "压缩摘要"]
    sx = Inches(0.5); sy = Inches(2.55); ssh = Inches(0.36); ggp = Inches(0.04)
    for i, n in enumerate(sources):
        y = sy + (ssh + ggp) * i
        # numbered
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, sx, y + Inches(0.04),
                               Inches(0.28), Inches(0.28))
        c.shadow.inherit = False; set_fill(c, ACCENT); set_line(c, None)
        tf = c.text_frame; tf.margin_left = tf.margin_right = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); set_run(r, str(i + 1), size=10, bold=True, color=WHITE)
        add_text(s, sx + Inches(0.4), y + Inches(0.05), Inches(4.5), Inches(0.3),
                 n, size=12.5, color=INK)

    # right: CLAUDE.md hierarchy
    rx = Inches(6.5)
    add_text(s, rx, Inches(2.0), Inches(6.0), Inches(0.4),
             "CLAUDE.md 层级 · 4 级", size=18, bold=True)

    headers = ["级别", "路径", "范围"]
    rows = [
        [[{"text": "托管", "bold": True}],
         [{"text": "/etc/claude-code/CLAUDE.md", "mono": True, "size": 10}],
         "系统范围（企业）"],
        [[{"text": "用户", "bold": True}],
         [{"text": "~/.claude/CLAUDE.md", "mono": True, "size": 10}],
         "用户级"],
        [[{"text": "项目", "bold": True}],
         [{"text": "CLAUDE.md", "mono": True, "size": 10},
          {"text": ".claude/rules/*.md", "mono": True, "size": 10, "newline": True}],
         "项目级"],
        [[{"text": "本地", "bold": True}],
         [{"text": "CLAUDE.local.md", "mono": True, "size": 10}],
         "个人 (gitignored)"],
    ]
    add_table(s, rx, Inches(2.55), Inches(6.3), Inches(2.1), headers, rows,
              body_size=10.5, header_size=11)
    tbl = s.shapes[-1].table
    tbl.columns[0].width = Inches(1.0)
    tbl.columns[1].width = Inches(3.0)
    tbl.columns[2].width = Inches(2.3)

    # memory card
    add_rect(s, rx, Inches(4.85), Inches(6.3), Inches(1.2),
             fill=WHITE, line=BORDER, line_pt=0.75, rounded=True, radius_ratio=0.06)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, rx, Inches(4.85),
                             Inches(0.08), Inches(1.2))
    bar.shadow.inherit = False; set_fill(bar, TEAL); set_line(bar, None)
    add_text(s, rx + Inches(0.2), Inches(4.95), Inches(6), Inches(0.3),
             "基于文件的记忆", size=11, bold=True, color=MUTED)
    add_text(s, rx + Inches(0.2), Inches(5.25), Inches(6), Inches(0.8),
             "不使用嵌入向量，无向量数据库。LLM 扫描文件头，按需挑出最多 5 个相关文件。"
             "完全可查看 / 可编辑 / 可纳入版本控制。",
             size=11, color=INK, line_spacing=1.55)

    # bottom design choice
    add_rect(s, Inches(0.5), Inches(6.25), Inches(12.3), Inches(0.85),
             fill=WHITE, line=BORDER, line_pt=0.75, rounded=True, radius_ratio=0.06)
    bar2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.5), Inches(6.25),
                              Inches(0.08), Inches(0.85))
    bar2.shadow.inherit = False; set_fill(bar2, ACCENT); set_line(bar2, None)
    add_text(s, Inches(0.7), Inches(6.32), Inches(12), Inches(0.3),
             "关键设计选择", size=11, bold=True, color=MUTED)
    add_multiline(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.5),
                  [{"text": "CLAUDE.md 作为", "size": 12},
                   {"text": "用户上下文", "size": 12, "bold": True},
                   {"text": "传递（模型遵从是", "size": 12},
                   {"text": "概率性", "size": 12, "bold": True},
                   {"text": "的），而非 system prompt（确定性）。真正提供确定性强制的，是",
                    "size": 12},
                   {"text": "权限规则", "size": 12, "bold": True},
                   {"text": "那一层。", "size": 12}], line_spacing=1.5)

    add_footer(s, "上下文与记忆", "11 / 15")


def slide_12_subagent(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(s)
    add_section_pill(s, Inches(0.5), Inches(0.45), "CHAPTER 08 · 子智能体")
    add_title(s, Inches(0.5), Inches(0.85),
              "子智能体委托 · SkillTool vs AgentTool", size=26)

    # 2 big cards
    cy = Inches(2.0); ch = Inches(1.7); cw = Inches(6.05)
    add_rect(s, Inches(0.5), cy, cw, ch, fill=WHITE, line=BORDER, line_pt=0.75,
             rounded=True, radius_ratio=0.06)
    bar1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), cy, cw, Inches(0.1))
    bar1.shadow.inherit = False; set_fill(bar1, TEAL); set_line(bar1, None)
    add_text(s, Inches(0.7), cy + Inches(0.2), cw - Inches(0.4), Inches(0.5),
             "SkillTool", size=20, bold=True, color=RGBColor(0x1F, 0x5E, 0x35))
    add_multiline(s, Inches(0.7), cy + Inches(0.7), cw - Inches(0.4), ch - Inches(0.85),
                  [{"text": "将", "size": 13},
                   {"text": "指令注入当前上下文", "size": 13, "bold": True},
                   {"text": "。", "size": 13},
                   {"text": "✓ 便宜，相同窗口", "size": 13, "newline": True},
                   {"text": "✓ 适合短小、可复用的能力片段", "size": 13, "newline": True}],
                  line_spacing=1.6)

    cx2 = Inches(0.5) + cw + Inches(0.2)
    add_rect(s, cx2, cy, cw, ch, fill=WHITE, line=BORDER, line_pt=0.75,
             rounded=True, radius_ratio=0.06)
    bar2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx2, cy, cw, Inches(0.1))
    bar2.shadow.inherit = False; set_fill(bar2, ACCENT); set_line(bar2, None)
    add_text(s, cx2 + Inches(0.2), cy + Inches(0.2), cw - Inches(0.4), Inches(0.5),
             "AgentTool", size=20, bold=True, color=BROWN_TXT)
    add_multiline(s, cx2 + Inches(0.2), cy + Inches(0.7), cw - Inches(0.4), ch - Inches(0.85),
                  [{"text": "生成", "size": 13},
                   {"text": "新的隔离上下文窗口", "size": 13, "bold": True},
                   {"text": "。", "size": 13},
                   {"text": "✗ 昂贵（约 7 倍 token）", "size": 13, "newline": True,
                    "color": RED_WARN},
                   {"text": "✓ 上下文安全，避免父级被污染", "size": 13, "newline": True}],
                  line_spacing=1.6)

    # 6 built-in types row
    add_text(s, Inches(0.5), Inches(3.95), Inches(12.3), Inches(0.4),
             "6 个内置类型 + 自定义", size=15, bold=True)
    types = ["Explore", "Plan", "General-purpose", "Claude Code Guide",
             "Verification", "Statusline-setup"]
    px = Inches(0.5); py = Inches(4.45); pw = Inches(1.85); pgp = Inches(0.08)
    for i, name in enumerate(types):
        add_pill(s, px + (pw + pgp) * i, py, pw, Inches(0.4), name,
                 fill=BLUE_BG, color=RGBColor(0x23, 0x4A, 0x6F), size=10.5)

    add_multiline(s, Inches(0.5), Inches(4.95), Inches(12.3), Inches(0.4),
                  [{"text": " + ", "size": 12, "color": MUTED},
                   {"text": ".claude/agents/*.md", "size": 12, "mono": True,
                    "color": ACCENT2},
                   {"text": " 自定义", "size": 12, "color": MUTED}])

    # 3 small cards
    cy3 = Inches(5.4); ch3 = Inches(1.65); cw3 = Inches(4.0); g3 = Inches(0.15)
    cards3 = [
        ("隔离模式（in-process 默认）",
         [("Worktree", "：Git worktree 文件系统隔离"),
          ("Remote", "：远程执行（内部专用）"),
          ("In-process", "：共享 FS，隔离对话")]),
        ("侧链转录稿",
         [("每个子智能体", "写自己的 .jsonl"),
          ("只把摘要", "回传父级"),
          ("POSIX flock()", " 协调，零外部依赖")]),
        ("权限覆盖",
         [("子智能体 permissionMode", " 生效；"),
          ("除非父级处于", " bypass / acceptEdits / auto"),
          ("显式用户决策", "始终优先")]),
    ]
    for i, (title, items) in enumerate(cards3):
        x = Inches(0.5) + (cw3 + g3) * i
        add_rect(s, x, cy3, cw3, ch3, fill=WHITE, line=BORDER, line_pt=0.75,
                 rounded=True, radius_ratio=0.06)
        add_text(s, x + Inches(0.2), cy3 + Inches(0.15), cw3 - Inches(0.4),
                 Inches(0.3), title, size=11, bold=True, color=MUTED)
        parts = []
        for j, (a, b) in enumerate(items):
            parts.append({"text": "• ", "size": 11, "color": ACCENT,
                          "newline": (j > 0)})
            parts.append({"text": a, "size": 11, "bold": True})
            parts.append({"text": b, "size": 11})
        add_multiline(s, x + Inches(0.2), cy3 + Inches(0.5),
                      cw3 - Inches(0.4), ch3 - Inches(0.65),
                      parts, line_spacing=1.5)

    add_footer(s, "子智能体委托", "12 / 15")


def slide_13_persistence(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(s)
    add_section_pill(s, Inches(0.5), Inches(0.45), "CHAPTER 09 · 持久化")
    add_title(s, Inches(0.5), Inches(0.85),
              "会话持久化 · 三个通道、仅追加 JSONL", size=26)

    headers = ["通道", "格式", "目的"]
    rows = [
        [[{"text": "会话转录稿", "bold": True}], "仅追加 JSONL · 链式修补",
         "完整对话；压缩边界记录 headUuid / anchorUuid / tailUuid"],
        [[{"text": "全局提示历史", "bold": True}],
         [{"text": "history.jsonl", "mono": True, "size": 11}],
         "跨会话提示召回（反向读取，用于上箭头）"],
        [[{"text": "子智能体侧链", "bold": True}], "每个子智能体独立 JSONL",
         "隔离子智能体历史；只摘要返回父级"],
    ]
    add_table(s, Inches(0.5), Inches(1.95), Inches(12.3), Inches(1.85),
              headers, rows, body_size=12)
    tbl = s.shapes[-1].table
    tbl.columns[0].width = Inches(2.6)
    tbl.columns[1].width = Inches(3.7)
    tbl.columns[2].width = Inches(6.0)

    # 2 cards
    cy = Inches(4.05); ch = Inches(1.35); cw = Inches(6.05)
    add_rect(s, Inches(0.5), cy, cw, ch, fill=WHITE, line=BORDER, line_pt=0.75,
             rounded=True, radius_ratio=0.06)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), cy, Inches(0.08), ch)
    bar.shadow.inherit = False; set_fill(bar, RED_WARN); set_line(bar, None)
    add_text(s, Inches(0.7), cy + Inches(0.15), cw - Inches(0.4), Inches(0.3),
             "恢复时权限永不自动恢复", size=11, bold=True, color=RED_WARN)
    add_text(s, Inches(0.7), cy + Inches(0.5), cw - Inches(0.4), ch - Inches(0.6),
             "信任总是在当前会话中重新建立。为守住这一安全不变量，"
             "系统宁可让用户多经历一次授权。",
             size=12, color=INK, line_spacing=1.5)

    cx2 = Inches(0.5) + cw + Inches(0.2)
    add_rect(s, cx2, cy, cw, ch, fill=WHITE, line=BORDER, line_pt=0.75,
             rounded=True, radius_ratio=0.06)
    bar2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx2, cy, Inches(0.08), ch)
    bar2.shadow.inherit = False; set_fill(bar2, ACCENT); set_line(bar2, None)
    add_text(s, cx2 + Inches(0.2), cy + Inches(0.15), cw - Inches(0.4), Inches(0.3),
             "设计权衡", size=11, bold=True, color=MUTED)
    add_multiline(s, cx2 + Inches(0.2), cy + Inches(0.5), cw - Inches(0.4),
                  ch - Inches(0.6),
                  [{"text": "仅追加 JSONL 体现一次取舍：", "size": 12},
                   {"text": "可审计性 + 简单性 > 查询能力", "size": 12, "bold": True},
                   {"text": "。每条事件可读、可版本控制，无需专用工具即可重建。",
                    "size": 12}], line_spacing=1.5)

    # bottom card
    by = Inches(5.55); bh = Inches(1.4)
    add_rect(s, Inches(0.5), by, Inches(12.3), bh, fill=LIGHT_BG,
             line=BORDER, line_pt=0.75, rounded=True, radius_ratio=0.04)
    add_text(s, Inches(0.7), by + Inches(0.15), Inches(12), Inches(0.3),
             "链式修补 + 文件历史检查点", size=11, bold=True, color=MUTED)
    add_multiline(s, Inches(0.7), by + Inches(0.5), Inches(12), bh - Inches(0.6),
                  [{"text": "会话加载器", "size": 12},
                   {"text": "读取时", "size": 12, "bold": True},
                   {"text": "修补消息链；磁盘上的数据不会被就地改写。", "size": 12,
                    "newline": False},
                   {"text": "--rewind-files", "mono": True, "size": 12,
                    "color": ACCENT2, "newline": True},
                   {"text": " 文件检查点存于 ", "size": 12},
                   {"text": "~/.claude/file-history/<sessionId>/", "mono": True,
                    "size": 12, "color": ACCENT2},
                   {"text": "。", "size": 12}], line_spacing=1.55)

    add_footer(s, "会话持久化", "13 / 15")


def slide_14_values(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(s)
    add_section_pill(s, Inches(0.5), Inches(0.45), "CHAPTER 10 · 设计哲学")
    add_title(s, Inches(0.5), Inches(0.85),
              "5 个价值观 → 13 条原则 → 实现", size=28)
    add_text(s, Inches(0.5), Inches(1.85), Inches(12.3), Inches(0.55),
             "架构里的每一条决策都能追溯到这 5 个人类价值观。",
             size=13, color=INK2)

    headers = ["价值观", "核心思想"]
    rows = [
        [[{"text": "人类决策权威", "bold": True}],
         [{"text": "人类通过主体层级保持控制；93% 提示批准率暴露批准疲劳后，应对是",
           "size": 12},
          {"text": "重新划分边界", "size": 12, "bold": True},
          {"text": "而非追加更多警告。", "size": 12}]],
        [[{"text": "安全 / 安保 / 隐私", "bold": True}],
         [{"text": "即使在人类警惕性下降时也能守住底线。", "size": 12},
          {"text": "7 个独立安全层。", "size": 12, "bold": True}]],
        [[{"text": "可靠执行", "bold": True}],
         [{"text": "按用户本意执行；", "size": 12},
          {"text": "收集—行动—验证", "size": 12, "bold": True},
          {"text": "闭环；优雅恢复。", "size": 12}]],
        [[{"text": "能力放大", "bold": True}],
         [{"text": "\"一个 Unix 工具，而不是产品。\" 98.4% 是让模型能工作的", "size": 12},
          {"text": "确定性基础设施", "size": 12, "bold": True},
          {"text": "。", "size": 12}]],
        [[{"text": "上下文适应性", "bold": True}],
         [{"text": "CLAUDE.md 层级、渐进式可扩展性，以及随时间演变的", "size": 12},
          {"text": "信任轨迹", "size": 12, "bold": True},
          {"text": "。", "size": 12}]],
    ]
    add_table(s, Inches(0.5), Inches(2.55), Inches(12.3), Inches(3.0),
              headers, rows, body_size=12)
    tbl = s.shapes[-1].table
    tbl.columns[0].width = Inches(3.0)
    tbl.columns[1].width = Inches(9.3)

    # 6th perspective card
    cy = Inches(5.75); ch = Inches(1.3)
    add_rect(s, Inches(0.5), cy, Inches(12.3), ch, fill=WHITE,
             line=BORDER, line_pt=0.75, rounded=True, radius_ratio=0.05)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), cy, Inches(0.08), ch)
    bar.shadow.inherit = False; set_fill(bar, TEAL); set_line(bar, None)
    add_text(s, Inches(0.7), cy + Inches(0.15), Inches(12), Inches(0.3),
             "第六个评估视角 · 长期能力保持", size=11, bold=True, color=MUTED)
    add_multiline(s, Inches(0.7), cy + Inches(0.5), Inches(12), ch - Inches(0.6),
                  [{"text": "论文援引证据：在 AI 辅助条件下工作的开发者，在", "size": 12},
                   {"text": "理解力测试中得分低 17%", "size": 12, "bold": True,
                    "color": RED_WARN},
                   {"text": "。架构既要放大短期能力，也不能侵蚀长期能力。", "size": 12}],
                  line_spacing=1.5)

    add_footer(s, "价值观与设计原则", "14 / 15")


def slide_15_summary(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.shadow.inherit = False
    set_fill(bg, DEEP_NAVY); set_line(bg, None)
    add_top_bar(s)

    add_text(s, Inches(0.5), Inches(0.85), Inches(8), Inches(0.4),
             "KEY TAKEAWAYS",
             size=12, bold=True, color=RGBColor(0xFF, 0xD9, 0xA8))

    add_text(s, Inches(0.5), Inches(1.4), Inches(12), Inches(1.2),
             "把投入放到 harness 上", size=46, bold=True, color=WHITE)

    add_multiline(s, Inches(0.5), Inches(2.85), Inches(12.3), Inches(1.8),
                  [{"text": "Claude Code 的核心洞察是：随着模型能力趋同，",
                    "size": 15, "color": RGBColor(0xD8, 0xDD, 0xE7)},
                   {"text": "harness 才是差异化的关键",
                    "size": 15, "bold": True,
                    "color": RGBColor(0xFF, 0xD9, 0xA8)},
                   {"text": "——钩子、分类器、压缩、隔离都不容易复制。",
                    "size": 15, "color": RGBColor(0xD8, 0xDD, 0xE7)},
                   {"text": "一个生产级编码智能体的工程性挑战，",
                    "size": 15, "color": RGBColor(0xD8, 0xDD, 0xE7), "newline": True},
                   {"text": "98.4% 在循环之外", "size": 15, "bold": True,
                    "color": RGBColor(0xFF, 0xD9, 0xA8)},
                   {"text": "。", "size": 15, "color": RGBColor(0xD8, 0xDD, 0xE7)}],
                  line_spacing=1.7)

    cards = [
        ("DESIGN", "为上下文稀缺、深度防御、可逆性加权风险而设计。"),
        ("SAFETY", "深度防御 + 拒绝优先；恢复会话时权限永不自动恢复。"),
        ("EXTENSIBILITY", "按上下文成本分层的扩展：Hooks / Skills / Plugins / MCP。"),
    ]
    cy = Inches(4.7); ch = Inches(1.5); cw = Inches(4.0); g = Inches(0.15)
    for i, (k, v) in enumerate(cards):
        x = Inches(0.5) + (cw + g) * i
        rect = add_rect(s, x, cy, cw, ch,
                        fill=RGBColor(0x2C, 0x37, 0x4E),
                        line=RGBColor(0x4A, 0x55, 0x6B), line_pt=0.75,
                        rounded=True, radius_ratio=0.06)
        add_text(s, x + Inches(0.25), cy + Inches(0.2), cw - Inches(0.5), Inches(0.4),
                 k, size=12, bold=True, color=RGBColor(0xFF, 0xD9, 0xA8))
        add_text(s, x + Inches(0.25), cy + Inches(0.6), cw - Inches(0.5),
                 ch - Inches(0.7), v, size=13,
                 color=RGBColor(0xE5, 0xE8, 0xEF), line_spacing=1.55)

    add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.3),
             "Source: Dive into Claude Code: The Design Space of Today's "
             "AI Agent System  ·  v2.1.88  ·  CC BY-NC-SA 4.0",
             size=10, italic=True, color=RGBColor(0xA8, 0xB0, 0xBD))

    add_text(s, Inches(11.5), Inches(7.15), Inches(1.5), Inches(0.3),
             "15 / 15", size=10, color=RGBColor(0xD8, 0xDD, 0xE7),
             align=PP_ALIGN.RIGHT, bold=True)


# =====================================================================
# Main
# =====================================================================

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_1_cover(prs)
    slide_2_tldr(prs)
    slide_3_four_questions(prs)
    slide_4_main_structure(prs)
    slide_5_layered(prs)
    slide_6_pipeline(prs)
    slide_7_compaction(prs)
    slide_8_defense(prs)
    slide_9_modes(prs)
    slide_10_extensibility(prs)
    slide_11_context_memory(prs)
    slide_12_subagent(prs)
    slide_13_persistence(prs)
    slide_14_values(prs)
    slide_15_summary(prs)

    out = os.path.join(os.path.dirname(__file__), "claudecode_ppt.pptx")
    prs.save(out)
    print(f"OK: wrote {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
